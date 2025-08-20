import os
from typing import Optional, Tuple
from pdf2docx import Converter
import pypandoc
import pandas as pd
from docx2pdf import convert as docx2pdf_convert
from PyPDF2 import PdfReader, PdfWriter
import glob
from concurrent.futures import ThreadPoolExecutor, as_completed
import configparser
from datetime import datetime
import argparse
import logging
from collections import deque
import threading
import time

# --- 日志配置 ---
# 集中式日志实例
logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)  # 设置默认日志级别

# 日志缓冲区和用于批量写入的锁
log_buffer = deque()
log_buffer_lock = threading.Lock()
LOG_BUFFER_SIZE = 10  # 每当日志条目达到10条或程序关闭时写入文件
LOG_WRITE_INTERVAL_SECONDS = 5  # 即使缓冲区未满，也每隔5秒写入文件


class DocumentConverter:

    def __init__(self, pdf_engine: str = "weasyprint", overwrite: bool = False, gui_mode: bool = False):
        """
        初始化文档转换器。
        加载配置文件，设置默认路径、PDF引擎、日志路径，并启动日志写入线程。
        :param pdf_engine: 用于PDF转换的默认引擎（如'weasyprint'）。
        :param overwrite: 是否覆盖现有输出文件。
        :param gui_mode: 是否在 GUI 模式下运行（概念性，影响线程行为）。
        """
        os.makedirs('./input', exist_ok=True)
        os.makedirs('./output', exist_ok=True)
        os.makedirs('./logs', exist_ok=True)
        self.config = configparser.ConfigParser()
        self._load_config()  # 自动加载配置文件

        # 从配置或传入参数获取设置，优先级：传入参数 > 配置文件 > 硬编码默认值
        self.pdf_engine = pdf_engine if pdf_engine != "weasyprint" else \
            self.config.get("Defaults", "default_pdf_engine", fallback="weasyprint")
        self.overwrite_output = overwrite if overwrite else \
            self.config.getboolean("Settings", "overwrite_output", fallback=False)
        self.gui_mode = gui_mode  # GUI 模式标志

        # 定义支持的转换格式映射
        self.supported_conversions = {
            'pdf': ['docx', 'txt', 'ppt'],
            'docx': ['pdf', 'md', 'html'],
            'xlsx': ['csv'],
            'md': ['docx', 'html'],
            'html': ['docx', 'pdf', 'md']
        }

        # 初始化日志文件路径
        self.log_path = self._generate_log_filename(
            self.config.get('Paths', 'log_directory', fallback='./logs')
        )

        # 设置文件日志处理器
        self._setup_file_logging()

        # 如果日志写入线程尚未启动，则启动它 (仅对第一个实例启动一次)
        if not hasattr(DocumentConverter, '_log_writer_thread_started'):
            DocumentConverter._log_writer_thread_started = True
            log_writer_thread = threading.Thread(target=self._run_log_writer, daemon=True)
            log_writer_thread.start()

    def _setup_file_logging(self):
        """为日志器设置文件处理器，特定于此实例的日志路径。"""
        # 移除任何现有文件处理器，以防止多次调用 __init__ 时重复日志记录
        for handler in logger.handlers[:]:
            if isinstance(handler, logging.FileHandler):
                logger.removeHandler(handler)
                handler.close()

        # 为当前日志文件添加新的文件处理器
        file_handler = logging.FileHandler(self.log_path, encoding='utf-8')
        formatter = logging.Formatter("[%(asctime)s] %(levelname)-8s: %(message)s")
        file_handler.setFormatter(formatter)
        logger.addHandler(file_handler)

    # 配置加载健壮性
    def _load_config(self, config_file: str = "config.ini"):
        """加载或创建配置文件 config.ini。"""
        # 设置默认配置
        default_config = {
            'Paths': {
                'input_directory': './input',
                'output_directory': './output',
                'log_directory': './logs'
            },
            'Defaults': {
                'default_pdf_engine': 'weasyprint',
                'default_docx_to_format': 'pdf',
                'default_pdf_to_format': 'docx',
                'default_general_to_format': 'pdf'
            },
            'Settings': {
                'overwrite_output': 'no',
                'pdf_chunk_size': '2'
            }
        }

        # 尝试加载配置文件
        if os.path.exists(config_file):
            try:
                # 优先尝试UTF-8编码
                self.config.read(config_file, encoding='utf-8')
                return  # 成功加载，直接返回
            except UnicodeDecodeError:
                try:
                    # 尝试其他常见编码
                    self.config.read(config_file, encoding='latin-1')
                    print(f"⚠️ 配置文件 {config_file} 使用 latin-1 编码加载。建议转换为 UTF-8 格式。")
                    return
                except Exception as e:
                    print(f"⚠️ 无法读取配置文件 {config_file}，将使用默认配置: {e}")
            except Exception as e:
                print(f"⚠️ 读取配置文件 {config_file} 时出错: {e}，将使用默认配置")

        # 创建默认配置
        self.config.read_dict(default_config)

        # 尝试写入默认配置
        try:
            with open(config_file, 'w', encoding='utf-8') as f:
                self.config.write(f)
            print(f"ℹ️ 已创建默认配置文件: {config_file}")
        except IOError as e:
            print(f"⚠️ 无法创建配置文件 {config_file}: {e}，将仅使用内存中的默认设置")


    def _generate_log_filename(self, log_dir: str) -> str:
        """生成并返回带有时间戳的日志文件路径。"""
        os.makedirs(log_dir, exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        return os.path.join(log_dir, f"conversion_log_{timestamp}.txt")

    def _add_log_to_buffer(self, level: str, file_path: str, message: str = ""):
        """将日志条目添加到缓冲区。"""
        now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        display_filepath = os.path.basename(file_path) if file_path and file_path != "N/A" else "N/A"
        log_entry = f"[{now}] {level:<8}: {display_filepath:<30} {message}"

        with log_buffer_lock:
            log_buffer.append(log_entry)
            # 当缓冲区达到大小限制时立即刷新
            if len(log_buffer) >= LOG_BUFFER_SIZE:
                self._flush_log_buffer()

    def _flush_log_buffer(self):
        """将缓冲的日志条目写入文件。"""
        with log_buffer_lock:
            if not log_buffer:
                return

            try:
                # 使用上下文管理器确保文件正确关闭
                with open(self.log_path, "a", encoding="utf-8") as log_file:
                    while log_buffer:
                        log_file.write(log_buffer.popleft() + "\n")
                # 清理后立即刷新IO缓冲区
                log_file.flush()
            except Exception as e:
                # 如果写入日志文件失败，则回退到打印
                print(f"⚠️ 写入缓冲日志到文件失败 ({self.log_path}): {e}")
                while log_buffer:
                    print(log_buffer.popleft())

    def _run_log_writer(self):
        """后台线程函数，用于定期刷新日志缓冲区。"""
        while True:
            time.sleep(LOG_WRITE_INTERVAL_SECONDS)
            with log_buffer_lock:
                if len(log_buffer) >= LOG_BUFFER_SIZE:
                    self._flush_log_buffer()
                # 即使缓冲区未满，也会根据间隔时间进行刷新

    # 修改 write_log 以使用缓冲区
    def write_log(self, status: str, file_path: str, message: str = ""):
        """
        向日志缓冲区写入消息。实际写入文件由后台线程处理。
        :param status: 操作状态 (如 "SUCCESS", "FAILED", "SKIPPED", "ERROR", "INFO")
        :param file_path: 关联的文件路径
        :param message: 详细消息
        """
        # 将自定义状态映射到日志级别
        log_level_map = {
            "SUCCESS": "INFO",
            "FAILED": "ERROR",
            "SKIPPED": "WARNING",
            "ERROR": "ERROR",
            "INFO": "INFO"
        }
        level = log_level_map.get(status.upper(), "INFO")

        # 使用标准日志模块进行即时控制台输出，并使用缓冲区进行文件写入
        msg = f"{os.path.basename(file_path) if file_path and file_path != 'N/A' else 'N/A'}: {message}"
        if level == "INFO":
            logger.info(msg)
        elif level == "WARNING":
            logger.warning(msg)
        elif level == "ERROR":
            logger.error(msg)

        self._add_log_to_buffer(status, file_path, message)

    def get_supported_conversions(self, from_format: str) -> list:
        """获取某格式支持转换的目标格式列表"""
        return self.supported_conversions.get(from_format.lower(), [])

    def convert(self, input_path: str, output_path: str,
                from_format: Optional[str] = None,
                to_format: Optional[str] = None) -> bool:
        """
        执行文档转换的主方法。
        :param input_path: 输入文件路径
        :param output_path: 输出文件路径
        :param from_format: 可选，手动指定输入文件格式
        :param to_format: 可选，手动指定输出文件格式
        :return: 转换是否成功
        """
        if not os.path.exists(input_path):
            self.write_log("ERROR", input_path, "输入文件不存在")
            print(f"❌ 错误: 输入文件不存在 - {input_path}")
            return False

        # 自动检测格式
        from_format = from_format or os.path.splitext(input_path)[1][1:].lower()
        to_format = to_format or os.path.splitext(output_path)[1][1:].lower()

        # 添加关键修复：防止将文件转换为相同格式
        if from_format == to_format:
            self.write_log("SKIPPED", input_path, f"源格式和目标格式相同 ({from_format})，无需转换")
            print(f"⚠️ 跳过: 源格式和目标格式相同 ({from_format})，无需转换 - {input_path}")
            return False

        # 检查是否支持该转换
        if to_format not in self.get_supported_conversions(from_format):
            self.write_log("ERROR", input_path, f"不支持从 {from_format} 转换为 {to_format}")
            print(f"❌ 错误: 不支持从 {from_format} 转换为 {to_format}")
            return False

        # 检查是否允许覆盖
        if os.path.exists(output_path) and not self.overwrite_output:
            self.write_log("SKIPPED", input_path, f"目标文件已存在，且不允许覆盖: {output_path}")
            print(f"⚠️ 跳过: 目标文件已存在，且不允许覆盖 - {output_path}")
            return False

        try:
            # 确保输出目录存在，否则可能导致转换失败
            os.makedirs(os.path.dirname(output_path), exist_ok=True)

            # 检查大文件并进行分块处理
            if from_format == 'pdf' and to_format == 'docx' and os.path.getsize(input_path) > 10 * 1024 * 1024:
                print(f"检测到大 PDF 文件 ({os.path.basename(input_path)})。正在分块转换...")
                self._pdf_to_docx_chunked(input_path, output_path)
            elif from_format == 'pdf' and to_format == 'txt' and os.path.getsize(input_path) > 10 * 1024 * 1024:
                print(f"检测到大 PDF 文件 ({os.path.basename(input_path)})。正在分块提取文本...")
                self._extract_text_from_pdf_chunked(input_path, output_path)
            elif from_format == 'pdf' and to_format == 'docx':
                self._pdf_to_docx(input_path, output_path)
            elif from_format == 'pdf' and to_format == 'txt':
                text = self.extract_text_from_pdf(input_path)
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(text)
            elif from_format == 'pdf' and to_format == 'pptx' and os.path.getsize(input_path) > 20 * 1024 * 1024:
                print(f"检测到大 PDF 文件 ({os.path.basename(input_path)})。正在分块转换为 PPT...")
                self._pdf_to_ppt_chunked(input_path, output_path)
            elif from_format == 'pdf' and to_format == 'pptx':
                self._pdf_to_ppt(input_path, output_path)
            elif from_format == 'docx' and to_format == 'pdf':
                # 统一 Word 转 PDF 逻辑，并在内部处理成功/失败及日志
                self._docx_to_pdf_robust(input_path, output_path)
            elif from_format == 'docx' and to_format in ['md', 'html']:
                self._docx_to_other(input_path, output_path, to_format)
            elif from_format == 'xlsx' and to_format == 'csv':
                self._excel_to_csv(input_path, output_path)
            elif from_format == 'md' and to_format == 'docx':
                self._markdown_to_docx(input_path, output_path)
            elif from_format == 'html' and to_format == 'pdf':
                self._html_to_pdf(input_path, output_path)
            else:
                # 使用 Pandoc 作为后备转换器
                self._convert_with_pandoc(input_path, output_path, from_format, to_format)

            self.write_log("SUCCESS", input_path, f"转换为 {output_path}")
            return True
        except Exception as e:
            err_msg = f"转换失败: {e}"
            print(f"❌ {err_msg} - 文件: {os.path.basename(input_path)}")
            self.write_log("FAILED", input_path, err_msg)
            return False

    def _pdf_to_docx(self, input_path: str, output_path: str):
        """内部方法：PDF转Word文档"""
        cv = Converter(input_path)
        cv.convert(output_path)
        cv.close()

    def _process_pdf_in_chunks(self, pdf_path: str, output_path: str,
                               process_chunk_func: callable,
                               chunk_size: int = None) -> None:
        """
        通用PDF分块处理函数
        :param pdf_path: PDF文件路径
        :param output_path: 输出路径
        :param process_chunk_func: 处理每个块的函数
        :param chunk_size: 分块大小，如果为None则从配置获取
        """
        reader = PdfReader(pdf_path)
        total_pages = len(reader.pages)

        # 获取配置的块大小，不同操作可能需要不同的默认值
        if chunk_size is None:
            chunk_size = self.config.getint('Settings', 'pdf_chunk_size', fallback=2)

        try:
            for i in range(0, total_pages, chunk_size):
                start_page = i
                end_page = min(i + chunk_size, total_pages)

                # 调用传入的处理函数
                process_chunk_func(reader, start_page, end_page, output_path)

                self.write_log("INFO", pdf_path, f"已处理 PDF 块 {start_page}-{end_page}")
        finally:
            # 确保reader资源被释放
            if hasattr(reader, '_file'):
                reader._file.close()

    def _pdf_to_docx_chunked(self, input_path: str, output_path: str):
        """分块转换大型 PDF 到 DOCX 以管理内存。"""
        temp_docx_files = []

        def process_chunk(reader, start_page, end_page, output_path):
            temp_pdf_path = f"{output_path}.temp_chunk_{start_page}-{end_page}.pdf"
            temp_docx_path = f"{output_path}.temp_chunk_{start_page}-{end_page}.docx"

            try:
                # 将块提取到临时 PDF
                writer = PdfWriter()
                for page_num in range(start_page, end_page):
                    writer.add_page(reader.pages[page_num])
                with open(temp_pdf_path, 'wb') as f:
                    writer.write(f)

                # 将临时 PDF 块转换为 DOCX
                cv = Converter(temp_pdf_path)
                cv.convert(temp_docx_path)
                cv.close()

                temp_docx_files.append(temp_docx_path)
            finally:
                # 确保临时文件被清理
                if os.path.exists(temp_pdf_path):
                    os.remove(temp_pdf_path)

        try:
            self._process_pdf_in_chunks(input_path, output_path, process_chunk)

            # 处理结果文件
            output_dir = os.path.splitext(output_path)[0] + "_chunks"
            os.makedirs(output_dir, exist_ok=True)
            for i, temp_file in enumerate(temp_docx_files):
                final_chunk_path = os.path.join(output_dir,
                                                f"{os.path.basename(output_path).replace('.docx', '')}_part{i + 1}.docx")
                os.rename(temp_file, final_chunk_path)
                self.write_log("SUCCESS", final_chunk_path, f"已生成 DOCX 块 {i + 1}")

            self.write_log("INFO", input_path,
                           f"分块 PDF 到 DOCX 转换完成。输出块在 '{output_dir}' 中。可能需要手动合并。")
            print(f"✅ 分块 PDF 到 DOCX 转换完成。输出块在 '{output_dir}' 中。您可能需要手动合并这些文件。")
        except Exception as e:
            # 清理所有临时文件
            for temp_file in temp_docx_files:
                if os.path.exists(temp_file):
                    try:
                        os.remove(temp_file)
                    except:
                        pass
            raise Exception(f"分块 PDF 到 DOCX 转换失败: {e}")

    def _docx_to_pdf_robust(self, input_path: str, output_path: str):
        """
        内部方法：鲁棒的 Word 转 PDF，优先使用 docx2pdf，失败则回退到 pypandoc。
        """
        # 尝试使用 docx2pdf
        try:
            output_dir = os.path.dirname(output_path)
            docx2pdf_convert(input_path, output_dir)

            # 处理docx2pdf生成的文件路径
            expected_generated_path = os.path.join(output_dir,
                                                   os.path.splitext(os.path.basename(input_path))[0] + ".pdf")
            if expected_generated_path != output_path and os.path.exists(expected_generated_path):
                os.rename(expected_generated_path, output_path)

            self.write_log("SUCCESS", input_path, "使用 docx2pdf 成功转换为 PDF")
            return
        except Exception as e:
            self.write_log("INFO", input_path, f"docx2pdf 失败，尝试回退到 pypandoc: {e}")
            print(f"⚠️ docx2pdf 转换失败 ({os.path.basename(input_path)})，尝试使用 pypandoc... 错误: {e}")

        # 回退到 pypandoc
        try:
            pypandoc.convert_file(
                input_path, "pdf",
                outputfile=output_path,
                extra_args=["--pdf-engine=" + self.pdf_engine]
            )
            self.write_log("SUCCESS", input_path, f"使用 pypandoc ({self.pdf_engine}) 成功转换为 PDF")
            return
        except Exception as e:
            error_msg = f"docx2pdf 和 pypandoc 均无法将 Word 转换为 PDF: {e}"
            self.write_log("FAILED", input_path, error_msg)
            raise Exception(error_msg)


    def _pdf_to_ppt(self, input_path: str, output_path: str):
        """内部方法：PDF转PPT（通过将每页PDF转换为图片并插入PPT）"""
        try:
            from pdf2image import convert_from_path
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError as e:
            raise ImportError("需要安装 pdf2image 和 python-pptx 库: pip install pdf2image python-pptx") from e

        # 创建新的PPT
        prs = Presentation()

        # 将PDF转换为图片
        images = convert_from_path(input_path)

        for i, image in enumerate(images):
            # 为每张图片创建新幻灯片
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)

            # 保存图片到临时文件
            temp_image_path = f"temp_page_{i}.jpg"
            image.save(temp_image_path, "JPEG")

            # 将图片添加到幻灯片
            left = top = Inches(0)
            slide.shapes.add_picture(temp_image_path, left, top, width=prs.slide_width, height=prs.slide_height)

            # 删除临时文件
            os.remove(temp_image_path)

        # 保存PPT
        prs.save(output_path)
        self.write_log("SUCCESS", input_path, f"PDF已成功转换为PPT: {output_path}")

    # 添加对大文件的PDF转PPT分块处理
    def _pdf_to_ppt_chunked(self, input_path: str, output_path: str, chunk_size: int = 5):
        """分块处理大型 PDF 转 PPT 以管理内存。"""
        try:
            from pdf2image import convert_from_path
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError as e:
            raise ImportError("需要安装 pdf2image 和 python-pptx 库: pip install pdf2image python-pptx") from e

        # 创建新的PPT
        prs = Presentation()

        # 读取PDF
        reader = PdfReader(input_path)
        total_pages = len(reader.pages)

        for i in range(0, total_pages, chunk_size):
            start_page = i
            end_page = min(i + chunk_size, total_pages)

            # 从PDF中提取页面范围
            temp_pdf_path = f"{output_path}.temp_chunk_{start_page}-{end_page}.pdf"
            writer = PdfWriter()
            for page_num in range(start_page, end_page):
                writer.add_page(reader.pages[page_num])
            with open(temp_pdf_path, 'wb') as f:
                writer.write(f)

            # 将PDF块转换为图片
            images = convert_from_path(temp_pdf_path)

            for image in images:
                # 为每张图片创建新幻灯片
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)

                # 保存图片到临时文件
                temp_image_path = f"temp_page_{i}.jpg"
                image.save(temp_image_path, "JPEG")

                # 将图片添加到幻灯片
                left = top = Inches(0)
                slide.shapes.add_picture(temp_image_path, left, top, width=prs.slide_width, height=prs.slide_height)

                # 删除临时文件
                os.remove(temp_image_path)

            # 清理临时PDF
            os.remove(temp_pdf_path)

            self.write_log("INFO", input_path, f"已处理 PDF 块 {start_page}-{end_page} 转 PPT")

        # 保存PPT
        prs.save(output_path)
        self.write_log("SUCCESS", input_path, f"分块PDF已成功转换为PPT: {output_path}")


    def _excel_to_csv(self, input_path: str, output_path: str):
        """内部方法：Excel转CSV"""
        data = pd.read_excel(input_path)
        data.to_csv(output_path, index=False)

    def _docx_to_other(self, input_path: str, output_path: str, to_format: str):
        """内部方法：Word转Markdown或HTML"""
        pypandoc.convert_file(input_path, to_format, outputfile=output_path)

    def _markdown_to_docx(self, input_path: str, output_path: str):
        """内部方法：Markdown转Word"""
        pypandoc.convert_file(input_path, 'docx', outputfile=output_path)

    def _html_to_pdf(self, input_path: str, output_path: str):
        """内部方法：HTML转PDF"""
        pypandoc.convert_file(
            input_path, 'pdf',
            outputfile=output_path,
            extra_args=self._get_pdf_engine_args()
        )

    def _convert_with_pandoc(self, input_path: str, output_path: str,
                             from_format: str, to_format: str):
        """使用 pandoc 进行通用格式转换 Markdown→PDF (或其他 Pandoc 支持的转换)"""
        extra_args = self._get_pdf_engine_args() if to_format == "pdf" else []
        pypandoc.convert_file(
            input_path, to_format,
            outputfile=output_path,
            extra_args=extra_args
        )

    def _get_pdf_engine_args(self) -> list:
        """根据当前设置返回 PDF 引擎参数"""
        return ["--pdf-engine=" + self.pdf_engine]

    def extract_text_from_pdf(self, pdf_path: str) -> str:
        """从PDF提取纯文本内容
        :param pdf_path: PDF文件路径
        :return: 提取的文本字符串
        """
        reader = PdfReader(pdf_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text

    def _extract_text_from_pdf_chunked(self, pdf_path: str, output_path: str):
        """分块从大型 PDF 中提取文本。"""
        reader = PdfReader(pdf_path)
        total_pages = len(reader.pages)
        chunk_size = self.config.getint('Settings', 'pdf_chunk_size', fallback=5)  # 文本提取的块大小

        all_text = []
        for i in range(0, total_pages, chunk_size):
            start_page = i
            end_page = min(i + chunk_size, total_pages)
            chunk_text = ""
            for page_num in range(start_page, end_page):
                chunk_text += reader.pages[page_num].extract_text() + "\n"
            all_text.append(chunk_text)
            self.write_log("INFO", pdf_path, f"已从 PDF 块 {start_page}-{end_page} 提取文本")

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write("".join(all_text))
        self.write_log("SUCCESS", output_path, "分块 PDF 文本提取完成")
        print(f"✅ 分块 PDF 文本提取完成。输出已保存到 '{output_path}'。")

    def merge_pdfs(self, pdf_paths: list, output_path: str):
        """合并多个PDF文件
        :param pdf_paths: PDF文件路径列表
        :param output_path: 合并后的输出路径
        """
        merger = PdfWriter()
        for pdf in pdf_paths:
            merger.append(pdf)
        merger.write(output_path)
        merger.close()

    def batch_convert(self, input_dir: str, output_dir: str,
                      from_format: str, to_format: str,
                      keep_structure: bool = False,
                      max_workers: int = 5) -> int:
        """
        批量转换目录下的文件
        """
        from_format = from_format.lower()
        to_format = to_format.lower()

        # 验证输入目录
        if not os.path.isdir(input_dir):
            self.write_log("ERROR", input_dir, "输入目录不存在或不是目录")
            print(f"❌ 错误: 输入目录不存在或不是目录 - {input_dir}")
            return 0

        # 验证转换支持
        if to_format not in self.get_supported_conversions(from_format):
            self.write_log("ERROR", f"{from_format} -> {to_format}", "不支持的批量转换类型")
            print(f"❌ 错误: 不支持从 {from_format} 批量转换为 {to_format}")
            return 0

        # 创建输出目录
        os.makedirs(output_dir, exist_ok=True)

        # 获取文件列表
        pattern = f"**/*.{from_format}"
        files = glob.glob(os.path.join(input_dir, pattern), recursive=True)

        if not files:
            self.write_log("INFO", input_dir, "未找到待转换文件")
            print(f"❗在 {input_dir} 中未找到 {from_format} 文件")
            return 0

        # 准备转换任务
        conversion_tasks = []
        for file in files:
            rel_path = os.path.relpath(file, input_dir)
            new_name = os.path.splitext(rel_path)[0] + f".{to_format}"

            if keep_structure:
                out_path = os.path.join(output_dir, new_name)
                os.makedirs(os.path.dirname(out_path), exist_ok=True)
            else:
                out_path = os.path.join(output_dir, os.path.basename(new_name))

            conversion_tasks.append((file, out_path))

        # 执行批量转换
        success_count = 0
        total_files = len(conversion_tasks)

        print(f"开始批量转换 {total_files} 个文件...")

        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            # 提交所有任务
            future_to_task = {
                executor.submit(self.convert, task[0], task[1], from_format, to_format): task
                for task in conversion_tasks
            }

            # 处理完成的任务
            for future in as_completed(future_to_task):
                task = future_to_task[future]
                try:
                    if future.result():
                        success_count += 1
                        print(f"✅ 成功: {os.path.basename(task[0])} -> {os.path.basename(task[1])}")
                    else:
                        print(f"❌ 失败: {os.path.basename(task[0])}")
                except Exception as e:
                    print(f"❌ 异常: {os.path.basename(task[0])} - {str(e)}")

        # 记录最终结果
        result_msg = f"批量转换完成。成功: {success_count}/{total_files}"
        self.write_log("INFO", input_dir, result_msg)
        print(f"{'✅' if success_count == total_files else '⚠️'} {result_msg}")

        return success_count


# --- main 函数保持不变，因为优化主要集中在类内部 ---
def main():
    parser = argparse.ArgumentParser(description="文档格式转换工具")
    parser.add_argument('input', help="输入文件或目录路径", nargs='?', default=None)
    parser.add_argument('output', help="输出文件或目录路径", nargs='?', default=None)
    parser.add_argument('--from', dest='from_format', help="源格式（必须在批量转换中指定）")
    parser.add_argument('--to', dest='to_format', help="目标格式")
    parser.add_argument('--batch', action='store_true', help="启用批量模式")
    # 添加一个可选参数用于覆盖现有文件
    parser.add_argument('--overwrite', action='store_true', help="如果输出文件已存在则覆盖")
    # 添加一个用于 GUI 模式的标志（概念性，影响线程行为）
    parser.add_argument('--gui', action='store_true', help="启用 GUI 模式（概念性，影响线程行为）")

    args = parser.parse_args()

    # 在实例化 DocumentConverter 时传递 overwrite 和 gui_mode 参数
    converter = DocumentConverter(overwrite=args.overwrite, gui_mode=args.gui)

    # 确定输入输出路径的最终值
    input_path = args.input or converter.config.get("Paths", "input_directory")
    output_path = args.output or converter.config.get("Paths", "output_directory")

    if args.batch:
        if not args.from_format or not args.to_format:
            print("❌ 批量模式下必须指定 --from 和 --to 参数")
            converter.write_log("ERROR", "N/A", "批量模式参数缺失")
            return

        print(
            f"🚀 开始批量转换 '{input_path}' 中的 {args.from_format} 文件到 '{output_path}' 为 {args.to_format} 格式...")
        count = converter.batch_convert(input_path, output_path,
                                        args.from_format, args.to_format)
        print(f"✅ 批量转换完成，共成功转换 {count} 个文件")
    else:
        # 如果没有指定 to_format，自动根据 config.ini 推断默认目标格式
        from_format = args.from_format
        to_format = args.to_format

        if from_format and not to_format:
            # 尝试根据源格式获取默认目标格式
            to_format = converter.config.get("Defaults", f"default_{from_format.lower()}_to_format", fallback=None)
            if not to_format:
                # 如果没有特定默认值，尝试获取通用默认值
                to_format = converter.config.get("Defaults", "default_general_to_format", fallback="pdf")

        # 如果最终 still 没有 to_format，且 input 也不是文件路径（无法推断），则报错
        if not to_format and not (input_path and os.path.isfile(input_path)):
            print("❌ 错误: 单文件转换时，如果未指定 --to 参数，且无法从输入文件推断，请手动指定目标格式。")
            converter.write_log("ERROR", input_path, "单文件转换目标格式缺失或无法推断")
            return
        elif not to_format and input_path and os.path.isfile(input_path):
            # 尝试从输出路径推断 to_format (如果用户只提供了输入，没有提供输出，但提供了输入文件)
            _, to_format_ext = os.path.splitext(output_path)  # 尝试从 output_path 推断，即使它是目录名
            to_format = to_format_ext[1:].lower() if to_format_ext else converter.config.get("Defaults",
                                                                                             "default_general_to_format",
                                                                                             fallback="pdf")
            print(f"ℹ️ 未指定目标格式，根据输出路径 '{output_path}' 或默认值推断为 '{to_format}'")

        print(f"🚀 开始转换 '{input_path}' ({from_format}) 到 '{output_path}' ({to_format})...")
        # 在 GUI 模式下，单文件转换也应该在后台线程中进行
        if args.gui:
            print("ℹ️ 在后台线程中运行单文件转换（GUI 模式概念性）。")
            with ThreadPoolExecutor(max_workers=1) as executor:
                future = executor.submit(converter.convert, input_path, output_path, from_format, to_format)
                success = future.result()  # 在此 CLI 上下文中等待完成
        else:
            success = converter.convert(input_path, output_path, from_format, to_format)

        if success:
            print(f"✅ 转换成功: {input_path} -> {output_path}")
        else:
            print(f"❌ 转换失败: {input_path} -> {output_path}")


if __name__ == '__main__':
    main()
    # 确保在退出前所有缓冲的日志都被写入
    with log_buffer_lock:
        if log_buffer:
            print("退出前刷新剩余日志...")
            # 为了能够访问 _flush_log_buffer 方法，需要一个 DocumentConverter 实例。
            # 如果 main() 没有运行或已完成，一个更健壮的解决方案可能涉及全局关闭钩子或传递转换器实例。
            # 在此上下文中，假设 main() 是主要入口，活跃转换器的刷新将得到处理。
            # 如果程序突然退出，一些日志可能会丢失。
            # 对于 `__main__` 来说，确保刷新的一个简单方法如下。
            try:
                # 临时创建一个转换器实例，仅用于访问 _flush_log_buffer
                # 这将重新初始化日志器，虽然不理想，但确保了 flush 方法的访问。
                temp_converter_for_flush = DocumentConverter()
                temp_converter_for_flush._flush_log_buffer()
            except Exception as e:
                print(f"最终日志刷新期间出错: {e}")